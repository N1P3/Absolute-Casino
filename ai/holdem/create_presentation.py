"""
Skrypt generujący prezentację PowerPoint dla projektu AI do pokera Texas Hold'em.
Przedmiot: Inteligencja Obliczeniowa
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Ścieżki do wykresów
BASE_PATH = r"C:\Users\avere\Desktop\studiaUG\projekt\absoluteCasino2\ai\holdem"
PLOTS_PATH = os.path.join(BASE_PATH, "runs", "run_20251114_125815", "plots")
VALIDATION_PATH = os.path.join(BASE_PATH, "validation_results")

def add_title_slide(prs, title, subtitle=""):
    """Dodaje slajd tytułowy."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Tło gradientowe (ciemne)
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(25, 25, 112)  # Midnight blue
    background.line.fill.background()

    # Tytuł
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # Podtytuł
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(200, 200, 200)
        p.alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(prs, title, bullet_points, image_path=None):
    """Dodaje slajd z treścią i opcjonalnie obrazkiem."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Nagłówek
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = RGBColor(25, 25, 112)
    header.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # Treść
    if image_path and os.path.exists(image_path):
        # Layout z obrazkiem po prawej
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5))
        slide.shapes.add_picture(image_path, Inches(5.2), Inches(1.5), width=Inches(4.5))
    else:
        # Pełna szerokość
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))

    tf = content_box.text_frame
    tf.word_wrap = True

    for i, point in enumerate(bullet_points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {point}"
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(50, 50, 50)
        p.space_after = Pt(12)

    return slide

def add_image_slide(prs, title, image_path, caption=""):
    """Dodaje slajd z dużym obrazkiem."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Nagłówek
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1))
    header.fill.solid()
    header.fill.fore_color.rgb = RGBColor(25, 25, 112)
    header.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # Obrazek
    if os.path.exists(image_path):
        slide.shapes.add_picture(image_path, Inches(0.8), Inches(1.2), width=Inches(8.4))

    # Podpis
    if caption:
        cap_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.5))
        tf = cap_box.text_frame
        p = tf.paragraphs[0]
        p.text = caption
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(100, 100, 100)
        p.alignment = PP_ALIGN.CENTER

    return slide

def add_two_column_slide(prs, title, left_content, right_content):
    """Dodaje slajd z dwoma kolumnami."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Nagłówek
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1))
    header.fill.solid()
    header.fill.fore_color.rgb = RGBColor(25, 25, 112)
    header.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # Lewa kolumna
    left_box = slide.shapes.add_textbox(Inches(0.3), Inches(1.3), Inches(4.5), Inches(5.5))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, point in enumerate(left_content):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = point
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(50, 50, 50)
        p.space_after = Pt(8)

    # Prawa kolumna
    right_box = slide.shapes.add_textbox(Inches(5), Inches(1.3), Inches(4.5), Inches(5.5))
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, point in enumerate(right_content):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = point
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(50, 50, 50)
        p.space_after = Pt(8)

    return slide

def add_architecture_slide(prs):
    """Dodaje slajd z architekturą modelu (ASCII art)."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Nagłówek
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.8))
    header.fill.solid()
    header.fill.fore_color.rgb = RGBColor(25, 25, 112)
    header.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Architektura modelu - PluribusPokerTransformer"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # Diagram
    arch_text = """
┌─────────────────────────────────────────────────────────────────┐
│                         INPUT LAYER                              │
├────────────────────────────┬────────────────────────────────────┤
│    Static State (18 dim)   │   Action Sequence (20 x 10 dim)   │
│  • Hole cards (2 karty)    │  • Historia akcji (20 akcji)       │
│  • Board cards (5 kart)    │  • Typ akcji (fold/call/raise)     │
│  • Stacks graczy (6)       │  • Kwota zakładu                   │
│  • Pot, Street             │  • ID gracza                       │
└────────────┬───────────────┴──────────────────┬─────────────────┘
             │                                   │
             ▼                                   ▼
┌────────────────────────┐         ┌────────────────────────────┐
│  STATIC TRANSFORMER    │         │    ACTION TRANSFORMER      │
│  • 2 warstwy           │         │    • 6 warstw              │
│  • 8 attention heads   │         │    • 8 attention heads     │
│  • Card embeddings     │         │    • Positional encoding   │
│  • CLS token           │         │    • CLS token             │
└────────────┬───────────┘         └──────────────┬─────────────┘
             │                                     │
             └──────────────┬──────────────────────┘
                            ▼
                ┌───────────────────────┐
                │     FUSION LAYER      │
                │  Concat → Linear →    │
                │  LayerNorm → GELU     │
                └───────────┬───────────┘
                     ┌──────┴──────┐
                     ▼             ▼
              ┌───────────┐  ┌───────────┐
              │ACTION HEAD│  │VALUE HEAD │
              │  3 klasy  │  │  Raise    │
              │F / C / R  │  │  amount   │
              └───────────┘  └───────────┘
"""

    arch_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.9), Inches(9.4), Inches(6.3))
    tf = arch_box.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = arch_text
    p.font.size = Pt(11)
    p.font.name = "Consolas"
    p.font.color.rgb = RGBColor(30, 30, 30)

    return slide

def create_presentation():
    """Główna funkcja tworząca prezentację."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slajd 1: Tytuł
    add_title_slide(prs,
        "AI do Texas Hold'em Poker",
        "Projekt z przedmiotu Inteligencja Obliczeniowa\nTransformer + Reinforcement Learning"
    )

    # Slajd 2: Agenda
    add_content_slide(prs, "Agenda", [
        "Wprowadzenie i cel projektu",
        "Architektura modelu (Transformer)",
        "Dataset i preprocessing",
        "Trening nadzorowany (Supervised Learning)",
        "Trening ze wzmocnieniem (Reinforcement Learning)",
        "Wyniki i ewaluacja",
        "Wnioski i podsumowanie"
    ])

    # Slajd 3: Wprowadzenie
    add_content_slide(prs, "Cel projektu", [
        "Stworzenie AI grającego w Texas Hold'em No-Limit",
        "Wykorzystanie architektury Transformer",
        "Dwuetapowy trening: SL + RL (PPO)",
        "Nauka z danych profesjonalnych graczy",
        "Optymalizacja poprzez self-play",
        "Eksport do formatu ONNX (deployment)"
    ])

    # Slajd 4: Problem pokerowy
    add_content_slide(prs, "Poker jako problem AI", [
        "Gra z niepełną informacją (ukryte karty)",
        "Duża przestrzeń stanów (~10^160 możliwych gier)",
        "Element losowości (karty, kolejność)",
        "Blef i psychologia (Game Theory Optimal)",
        "Wymagana analiza sekwencji akcji",
        "Multi-task: wybór akcji + kwota raise'a"
    ])

    # Slajd 5: Architektura
    add_architecture_slide(prs)

    # Slajd 6: Parametry modelu
    add_two_column_slide(prs, "Parametry modelu",
        [
            "ARCHITEKTURA:",
            "• d_model: 512",
            "• Attention heads: 8",
            "• Action Transformer: 6 warstw",
            "• Static Transformer: 2 warstwy",
            "• Feed-forward: 2048",
            "• Dropout: 0.1",
            "",
            "EMBEDDINGS:",
            "• Rank embedding: 13 → 8 dim",
            "• Suit embedding: 4 → 8 dim",
            "• Position embedding: 7 pozycji",
            "• Stage embedding: 4 ulice"
        ],
        [
            "ROZMIAR:",
            "• ~28 milionów parametrów",
            "• ~112 MB (fp32)",
            "• ~56 MB (fp16)",
            "",
            "WYJŚCIA:",
            "• Action logits: [FOLD, CALL, RAISE]",
            "• Value head: log(raise_multiplier)",
            "",
            "INNOWACJE:",
            "• Pre-norm architecture",
            "• Learnable CLS tokens",
            "• Unknown card masking"
        ]
    )

    # Slajd 7: Dataset
    add_content_slide(prs, "Dataset - źródła danych", [
        "ACPC 2017 - Annual Computer Poker Competition",
        "  Gry bot vs bot z oficjalnych zawodów",
        "HandHQ - profesjonalne rozgrywki heads-up",
        "  Prawdziwe gry zawodowych graczy",
        "Pluribus (Facebook AI Research)",
        "  Gry AI na poziomie mistrzowskim",
        "Łącznie: setki tysięcy punktów decyzyjnych",
        "Format: NPZ (NumPy compressed) - 55.5 MB"
    ])

    # Slajd 8: Encoding danych
    add_two_column_slide(prs, "Kodowanie danych",
        [
            "STATIC STATE (18 cech):",
            "• Hole cards: 2 ID (0-52)",
            "• Board cards: 5 ID (0-52)",
            "• Stacks: 6 znormalizowanych",
            "• Pot: 1 znormalizowany",
            "• Street: 4 one-hot",
            "  (preflop/flop/turn/river)",
            "",
            "KODOWANIE KART:",
            "card_id = rank * 4 + suit",
            "0-51: karty, 52: nieznana"
        ],
        [
            "ACTION SEQUENCE (20 x 10):",
            "• Player ID: 6 one-hot",
            "• Action type: 3 one-hot",
            "  (fold / call / raise)",
            "• Bet amount: 1 float",
            "",
            "PREPROCESSING:",
            "• Rust (10-100x szybszy)",
            "• Parsowanie TOML/PHH",
            "• Symulacja gry",
            "• Ekstrakcja punktów decyzji"
        ]
    )

    # Slajd 9: Trening SL
    add_content_slide(prs, "Trening nadzorowany (Supervised Learning)", [
        "Optimizer: AdamW (lr=1e-4, weight_decay=0.01)",
        "Batch size: 1024",
        "Epoki: 20",
        "Mixed precision (fp16) - 2x szybszy trening",
        "Gradient clipping: max_norm=1.0",
        "Scheduler: ReduceLROnPlateau (factor=0.5)",
        "Outcome-weighted loss - nauka z wygranych"
    ], os.path.join(PLOTS_PATH, "train_epoch_accuracy.png"))

    # Slajd 10: Funkcja straty
    add_content_slide(prs, "Multi-task Loss Function", [
        "L_action = CrossEntropyLoss(action_logits, target)",
        "  Klasyfikacja: FOLD / CALL / RAISE",
        "",
        "L_value = MSELoss(log_pred, log_target)",
        "  Regresja kwoty raise'a (skala log)",
        "",
        "L_total = 1.0 * L_action + 0.2 * L_value",
        "",
        "Outcome weighting:",
        "  weight = 0.2 + 1.8 * sigmoid(outcome / T)",
        "  Zakres wag: [0.2, 2.0]"
    ])

    # Slajd 11: Krzywe treningowe
    add_image_slide(prs, "Krzywe treningowe - Loss",
        os.path.join(PLOTS_PATH, "train_epoch_loss_total.png"),
        "Spadek funkcji straty podczas 20 epok treningu nadzorowanego"
    )

    # Slajd 12: RL Training
    add_content_slide(prs, "Reinforcement Learning - PPO", [
        "Algorytm: Proximal Policy Optimization (PPO)",
        "Self-play w środowisku PokerKit",
        "Epizody: 50,000 rozdań",
        "Learning rate: 2e-5 (niższy niż SL)",
        "Clip epsilon: 0.25",
        "Entropy coefficient: 0.02",
        "Opponent pool: 5 modeli (85%) + random (15%)",
        "Target KL: 0.02 (early stopping)"
    ])

    # Slajd 13: Środowisko RL
    add_content_slide(prs, "Środowisko treningowe RL", [
        "PokerKit NoLimitTexasHoldem",
        "Gracze: 2 (heads-up)",
        "Starting stack: 1000 chipów",
        "Blinds: 5/10",
        "Randomized stacks: Tak",
        "",
        "Opponent Pool Strategy:",
        "  - Zapobiega overfittingowi do jednego przeciwnika",
        "  - Aktualizacja co 500 epizodów",
        "  - 15% gier vs random baseline"
    ])

    # Slajd 14: Confusion Matrix
    add_image_slide(prs, "Confusion Matrix - Walidacja",
        os.path.join(VALIDATION_PATH, "confusion_matrix_Test.png"),
        "Macierz pomyłek dla klasyfikacji akcji (FOLD/CALL/RAISE)"
    )

    # Slajd 15: Metryki
    add_image_slide(prs, "Metryki wydajności",
        os.path.join(VALIDATION_PATH, "metrics_summary_Test.png"),
        "Per-class accuracy, precision, recall i F1-score"
    )

    # Slajd 16: Porównanie modeli
    add_image_slide(prs, "Porównanie: SL vs RL Model",
        os.path.join(BASE_PATH, "profit_graph.png"),
        "Kumulatywny profit w 500 grach - model RL znacząco lepszy"
    )

    # Slajd 17: Wyniki porównania
    add_two_column_slide(prs, "Wyniki: SL vs RL (100 gier)",
        [
            "SUPERVISED MODEL:",
            "• Win rate: 26%",
            "• Profit: +1.84 BB",
            "",
            "Akcje:",
            "• FOLD: 49",
            "• CALL: 310",
            "• RAISE: 5",
            "",
            "Charakter: konserwatywny",
            "zbyt dużo fold/call"
        ],
        [
            "RL MODEL:",
            "• Win rate: 74%",
            "• Poprawa: +185%",
            "",
            "Akcje:",
            "• FOLD: 0",
            "• CALL: 2",
            "• RAISE: 233",
            "",
            "Charakter: agresywny",
            "dominuje poprzez raise'y"
        ]
    )

    # Slajd 18: Wins comparison
    add_image_slide(prs, "Liczba wygranych",
        os.path.join(BASE_PATH, "rl_wins_by_player.png"),
        "Rozkład wygranych: Supervised (~115), RL (~172), Heuristic (~208)"
    )

    # Slajd 19: Deployment
    add_content_slide(prs, "Deployment i eksport", [
        "Format eksportu: ONNX (Open Neural Network Exchange)",
        "Rozmiar: 895 KB + 104 MB (external weights)",
        "Cross-platform: C++, Java, C#, JavaScript",
        "Bez zależności od PyTorch",
        "",
        "Inference engine:",
        "  - Batch processing",
        "  - Interactive play mode",
        "  - Temperature sampling"
    ])

    # Slajd 20: Podsumowanie
    add_content_slide(prs, "Podsumowanie", [
        "Stworzono działającego agenta AI do Texas Hold'em",
        "Architektura Transformer z dual-encoder design",
        "Trening na danych profesjonalnych + self-play RL",
        "Model RL wygrywa 74% gier vs model SL (26%)",
        "Agresywna strategia (raise) dominuje nad pasywną",
        "Eksport ONNX umożliwia deployment produkcyjny",
        "",
        "Technologie: PyTorch, PPO, PokerKit, Rust, ONNX"
    ])

    # Slajd 21: Pytania
    add_title_slide(prs,
        "Pytania?",
        "Dziękuję za uwagę"
    )

    # Zapisz prezentację
    output_path = os.path.join(BASE_PATH, "Prezentacja_Poker_AI.pptx")
    prs.save(output_path)
    print(f"Prezentacja zapisana: {output_path}")
    return output_path

if __name__ == "__main__":
    create_presentation()
